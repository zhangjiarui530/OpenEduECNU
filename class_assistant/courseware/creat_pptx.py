import streamlit as st
import os
import time
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
from localApp import generate_text_from_model
from extract_high_freq_words import extract_high_freq_words_from_file

# 创建空白演示文稿
prs = Presentation()


# 解析文本，将每一页的内容保存到一个字典中。
def parse_text_to_pages(text):
    pages_content = {}  # 创建一个空字典用于存储每页的内容

    pages = text.split('第')[1:]  # 跳过第一个空字符串

    for page in pages:
        try:
            # 尝试拆分页面内容
            page_number, page_content = page.split('页：', 1)
            page_number = page_number.strip()  # 清除页码两端可能存在的空白字符

        except ValueError as e:
            # 如果拆分失败，打印错误消息并跳过当前循环
            print(f"Error: {e}. Please check the format of the text around '第{page}'.")
            continue

        # 按段落拆分页面内容并去除两端的空白字符
        paragraphs = [paragraph.strip() for paragraph in page_content.split('\n') if paragraph.strip()]

        # 将页面内容（不包括空行）保存到字典中
        pages_content[page_number] = '\n'.join(paragraphs)

    return pages_content


# 遍历字典中的每一页内容，并根据内容的前缀执行打印操作。
def process_pages(pages_dict):
    for page_number, text in pages_dict.items():

        print(f"处理第{page_number}页内容：")

        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)

        left = top = Inches(0)  # 图像的起始位置
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # 将图像添加为幻灯片上的形状，覆盖整个幻灯片
        pic = slide.shapes.add_picture(r"C:\开放原子\4.17\class_assistant\pic\3.png", left, top, slide_width,
                                       slide_height)

        for line in text.strip().split("\n"):

            if line.startswith("一级标题："):

                print("大标题:", line.replace("一级标题：", ""))
                title_shape = slide.shapes.title
                title_shape.text = line.replace("一级标题：", "")

                print("一级标题:", line.replace("一级标题：", ""))
                content_shape = slide.placeholders[1]
                content_shape.text = line.replace("一级标题：", "")


            elif line.startswith("二级标题"):
                # 一次性替换所有二级标题的编号
                line = line.replace("二级标题", "")
                for i in range(1, 4):
                    line = line.replace(f"{i}：", "")
                print("二级标题:", line)
                p = content_shape.text_frame.add_paragraph()
                p.text = line
                p.level = 1  # 增加缩进级别


            elif line.startswith("正文："):
                print("正文:", line.replace("正文：", ""))
                p = content_shape.text_frame.add_paragraph()
                p.text = line.replace("正文：", "")
                p.level = 2
                p.font.size = Pt(20)

            else:
                continue

        print()  # 在每页内容之后打印一个空行以分隔页面


def create_pptx(course_name, unit_name, course_num, course_time, course_sub, file):
    st.warning('提交成功，教学PPT生成中，请稍后。。。。。。')

    # 创建一个placeholder，稍后用于显示或隐藏加载圈
    placeholder = st.empty()

    # 使用Streamlit的markdown和HTML功能来创建一个转动的加载圈
    # 这里使用的是简单的HTML和CSS
    placeholder.markdown("""
        <div style="position: fixed; top: 50%; left: 50%; transform: translate(-50%, -50%);">
            <div style="border: 20px solid #f3f3f3; border-top: 20px solid #3498db; border-radius: 50%; width: 160px; height: 160px; animation: spin 4s linear infinite;">
            </div>
        </div>
        <style>
        @keyframes spin {
            0% { transform: rotate(0deg); }
            100% { transform: rotate(360deg); }
        }
        </style>
    """, unsafe_allow_html=True)

    high_freq_words = extract_high_freq_words_from_file(file, 10)

    prompt1 = "现在你设计一份" + course_name + "课程其中的一节课的PPT大纲，课程内容为：" + unit_name + ",适用专业为" + course_sub + ",本课程" + course_name + f"中频繁出现的关键词包括{high_freq_words}。" + """
    设计六页左右，每页设计内容包含PPT的页面标题、一级标题、二级标题和正文，
    同一个一级标题下需要有3个或者2个二级标题和这些二级标题所对应的正文，
    每个二级标题下必须要有一个正文，
    需分别注明是一级标题、二级标题或正文
    一级标题不超过15个中文字，
    二级标题不超过15个中文字，
    正文不少于100个中文字

    输出格式如下，每一页都需输出：
    第X页：XXX
    一级标题：XXX
    二级标题：XXX
    三级标题：XXX
    正文：XXX
    """

    text = ""
    text = generate_text_from_model(prompt1)

    # 添加标题布局的幻灯片
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    left = top = Inches(0)  # 图像的起始位置
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 将图像添加为幻灯片上的形状，覆盖整个幻灯片
    pic = slide.shapes.add_picture(r"C:\开放原子\4.17\class_assistant\pic\2.png", left, top, slide_width, slide_height)

    # 设置标题和副标题
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = course_name
    subtitle.text = course_num + "    " + unit_name

    pages_dict = parse_text_to_pages(text)
    process_pages(pages_dict)

    # 使用placeholder.empty()来隐藏加载圈
    placeholder.empty()

    # 保存文档
    prs.save('result\\' + course_name + '课程' + unit_name + '课时教学PPT.pptx')

    st.success('教学PPT生成成功！')

    time.sleep(1)

    os.startfile('result')
    os.startfile('result\\' + course_name + '课程' + unit_name + '课时教学PPT.pptx')
